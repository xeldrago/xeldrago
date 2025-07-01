from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

### ---- SLIDE 1: Title Slide ---- ###
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Sanctioned vs Working Strength"
subtitle.text = "Chennai GST & Customs - April 2025"

title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 102, 153)

### ---- SLIDE 2: Table Slide ---- ###
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Strength Overview"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.bold = True

rows, cols = 4, 4
left = Inches(1)
top = Inches(1.5)
width = Inches(7)
height = Inches(1.5)

table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set header
headers = ["Zone", "Sanctioned", "Working", "Vacancy"]
data = [
    ["Chennai GST", 177, 123, 54],
    ["Chennai Customs", 101, 71, 30]
]

for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    para = cell.text_frame.paragraphs[0]
    para.font.bold = True
    para.font.size = Pt(16)
    para.font.color.rgb = RGBColor(255, 255, 255)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 75, 150)

# Add data
for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = str(value)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        cell.fill.solid()
        if row_idx % 2 == 1:
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        else:
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

### ---- SLIDE 3: Text Visualization ---- ###
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Vacancy Visualization"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

textbox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(6), Inches(1))
frame = textbox.text_frame
frame.text = "This chart shows the current vacancy across zones."
para = frame.paragraphs[0]
para.font.size = Pt(14)
para.font.italic = True
para.font.color.rgb = RGBColor(89, 89, 89)

### ---- Save the presentation ---- ###
prs.save("Sanctioned_vs_Working_Strength.pptx")
print("Presentation created successfully.")
